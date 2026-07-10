"""Build the redesigned DataVerse AI FYP defense deck (accurate + verifiability-first)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----
INK      = RGBColor(0x0F, 0x17, 0x2A)
MUTED    = RGBColor(0x47, 0x55, 0x69)
SLATE    = RGBColor(0x64, 0x74, 0x8B)
VIOLET   = RGBColor(0x7C, 0x3A, 0xED)
INDIGO   = RGBColor(0x31, 0x2E, 0x81)
DEEP     = RGBColor(0x1E, 0x1B, 0x4B)
EMERALD  = RGBColor(0x10, 0xB9, 0x81)
LIGHT    = RGBColor(0xF8, 0xFA, 0xFC)
BORDER   = RGBColor(0xE2, 0xE8, 0xF0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ICE      = RGBColor(0xC7, 0xD2, 0xFE)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
CARDBG   = RGBColor(0xF5, 0xF3, 0xFF)

HEAD = "Cambria"
BODY = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return r


def rrect(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.08):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def circle(s, x, y, d, fill, ch="", txtcolor=WHITE, size=16):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = fill; c.line.fill.background(); c.shadow.inherit = False
    if ch:
        tf = c.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ch
        r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = txtcolor; r.font.name = BODY
    return c


def text(s, t, x, y, w, h, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT,
         font=BODY, anchor=MSO_ANCHOR.TOP, italic=False, line_spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04); tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = t.split("\n") if isinstance(t, str) else t
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = font
    return tb


def bullets(s, items, x, y, w, h, size=15, color=MUTED, gap=6, bold_lead=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04); tf.margin_top = Inches(0.02)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.05
        lead, rest = (it if isinstance(it, tuple) else (None, it))
        r = p.add_run(); r.text = "•  "; r.font.size = Pt(size); r.font.color.rgb = VIOLET; r.font.name = BODY
        if lead:
            rb = p.add_run(); rb.text = lead + "  "; rb.font.size = Pt(size); rb.font.bold = True; rb.font.color.rgb = INK; rb.font.name = BODY
        rr = p.add_run(); rr.text = rest; rr.font.size = Pt(size); rr.font.color.rgb = color; rr.font.name = BODY
    return tb


def kicker(s, t, x=0.7, y=0.55, color=VIOLET):
    text(s, t.upper(), x, y, 9, 0.35, size=12.5, color=color, bold=True, font=BODY)


def title(s, t, x=0.7, y=0.92, w=11.9, color=INK, size=33):
    text(s, t, x, y, w, 0.9, size=size, color=color, bold=True, font=HEAD)


def pagenum(s, n):
    text(s, str(n), 12.7, 7.04, 0.5, 0.3, size=11, color=SLATE, align=PP_ALIGN.RIGHT)
    text(s, "DataVerse AI", 0.7, 7.04, 4, 0.3, size=10, color=SLATE)


def badge(s, t, x, y, fill=EMERALD, w=1.15, txtcolor=WHITE):
    b = rrect(s, x, y, w, 0.32, fill=fill, radius=0.5)
    tf = b.text_frame; tf.margin_top = tf.margin_bottom = 0; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = t; r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = txtcolor; r.font.name = BODY
    return b


def card(s, x, y, w, h, num, head, body, accent=VIOLET, headsize=15, bodysize=12.5, unique=False):
    rrect(s, x, y, w, h, fill=WHITE, line=BORDER, line_w=1.0, radius=0.06)
    circle(s, x + 0.22, y + 0.22, 0.5, accent, num, WHITE, 15)
    text(s, head, x + 0.85, y + 0.2, w - 1.0, 0.5, size=headsize, color=INK, bold=True, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
    text(s, body, x + 0.24, y + 0.78, w - 0.48, h - 0.95, size=bodysize, color=MUTED, font=BODY, line_spacing=1.04)
    if unique:
        badge(s, "★ UNIQUE", x + w - 1.15, y + 0.16, EMERALD, 1.0)


# ============================================================ SLIDE 1 — TITLE
s = slide(); bg(s, DEEP)
circle(s, 0.7, 0.7, 0.62, VIOLET, "DV", WHITE, 19)
text(s, "DataVerse AI", 1.5, 0.72, 8, 0.6, size=20, color=WHITE, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
text(s, "Computing & Technology Department  ·  Iqra University, Islamabad (H9 Campus)", 0.7, 6.95, 12, 0.4, size=12, color=ICE)
text(s, "DataVerse AI", 0.7, 2.45, 12, 1.0, size=58, color=WHITE, bold=True, font=HEAD)
text(s, "The Verifiable AI Data Analyst", 0.72, 3.5, 12, 0.6, size=26, color=ICE, font=HEAD, italic=True)
text(s, "A deterministic, two-agent business-intelligence platform where every number\ncan be proven, re-verified, and trusted — not just generated.",
     0.72, 4.2, 11.5, 0.9, size=15, color=RGBColor(0xCB,0xD5,0xE1), line_spacing=1.15)
rrect(s, 0.7, 5.35, 7.7, 1.15, fill=RGBColor(0x2A,0x26,0x5F), radius=0.08)
text(s, "Presented by", 0.95, 5.5, 4, 0.3, size=11, color=ICE, bold=True)
text(s, "Muhammad Muavia (37766)    ·    Hasnain Ali (37721)    ·    M. Shahan Subhani (37655)",
     0.95, 5.82, 7.3, 0.5, size=13.5, color=WHITE, bold=True)
text(s, "Supervisor: Ms. Tayyaba Shehzad        BS Computer Science  ·  2025–26",
     0.95, 6.15, 7.3, 0.3, size=12, color=ICE)

# ============================================================ SLIDE 2 — AGENDA
s = slide(); bg(s, WHITE); kicker(s, "Agenda"); title(s, "What we'll cover")
items_l = [("01", "Introduction & Problem"), ("02", "Proposed Solution"),
           ("03", "Two-Agent Architecture"), ("04", "The Verifiability Stack"),
           ("05", "Unique Features")]
items_r = [("06", "Technology Stack"), ("07", "Dataset & Testing"),
           ("08", "Results & Verification"), ("09", "How We Compare"),
           ("10", "Conclusion & Future Work")]
for col, items, x in ((0, items_l, 0.9), (1, items_r, 7.0)):
    for i, (n, t) in enumerate(items):
        y = 1.95 + i * 0.92
        circle(s, x, y, 0.56, VIOLET if col == 0 else INDIGO, n, WHITE, 14)
        text(s, t, x + 0.78, y + 0.02, 5.0, 0.55, size=17, color=INK, bold=True, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
pagenum(s, 2)

# ============================================================ SLIDE 3 — INTRODUCTION
s = slide(); bg(s, WHITE); kicker(s, "01  ·  Introduction"); title(s, "Data is everywhere — insight is not")
bullets(s, [
    ("Data overload.", "Businesses generate huge volumes of sales, customer, and operational data every day."),
    ("A skills gap.", "Turning raw data into insight needs Python, SQL, statistics, and ML — skills most decision-makers don't have."),
    ("Slow decisions.", "Non-technical staff wait on analysts; static dashboards and opaque models add friction."),
    ("The LLM opportunity — and its flaw.", "LLMs let anyone 'talk to data', but they can hallucinate numbers, so their output can't be trusted for real decisions."),
], 0.9, 1.95, 7.0, 4.5, size=15.5)
rrect(s, 8.25, 1.95, 4.35, 4.3, fill=DEEP, radius=0.06)
text(s, "The core question", 8.55, 2.2, 3.8, 0.4, size=13, color=ICE, bold=True)
text(s, "“If an AI gives me a number,\nhow do I know it's right?”", 8.55, 2.7, 3.8, 1.4, size=21, color=WHITE, bold=True, font=HEAD, line_spacing=1.1)
text(s, "DataVerse AI answers this by computing every number deterministically — and letting you prove it.",
     8.55, 4.35, 3.8, 1.6, size=14, color=ICE, line_spacing=1.2)
pagenum(s, 3)

# ============================================================ SLIDE 4 — PROBLEM STATEMENT
s = slide(); bg(s, WHITE); kicker(s, "02  ·  Problem Statement"); title(s, "Why existing tools fall short")
cols = [
    ("Traditional BI", "Tableau / Power BI", ["Needs analysts & setup", "Static dashboards", "No natural-language Q&A", "No per-number explanation"], SLATE),
    ("LLM analysts", "ChatGPT ADA / Julius", ["Conversational & fast", "BUT numbers can be hallucinated", "Outputs not reproducible", "Nothing to audit or re-verify"], AMBER),
    ("The gap", "What's missing", ["Conversational AND trustworthy", "Deterministic numbers", "Provable, re-verifiable results", "Usable by non-technical staff"], EMERALD),
]
for i, (h, sub, its, ac) in enumerate(cols):
    x = 0.9 + i * 4.05
    rrect(s, x, 1.95, 3.75, 4.35, fill=WHITE, line=BORDER, line_w=1.2, radius=0.05)
    rrect(s, x, 1.95, 3.75, 0.75, fill=ac, radius=0.05)
    text(s, h, x + 0.25, 2.0, 3.3, 0.4, size=16, color=WHITE, bold=True, font=BODY)
    text(s, sub, x + 0.25, 2.36, 3.3, 0.3, size=11.5, color=WHITE, font=BODY)
    bullets(s, its, x + 0.18, 2.95, 3.45, 3.2, size=13)
text(s, "Problem: there is no conversational BI platform that is both easy for non-technical users and provably trustworthy.",
     0.9, 6.55, 11.6, 0.6, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER, font=BODY)
pagenum(s, 4)

# ============================================================ SLIDE 5 — PROPOSED SOLUTION
s = slide(); bg(s, WHITE); kicker(s, "03  ·  Proposed Solution"); title(s, "DataVerse AI — chat with your data, then prove it")
text(s, "Upload a CSV/Excel file and ask questions in plain English. Two specialised agents profile the data, compute "
        "business metrics, EDA, and an optional prediction — all deterministically in Pandas / scikit-learn. The LLM only "
        "polishes wording; it never invents a number. With no API keys, everything still runs in offline “Mock mode.”",
     0.9, 1.85, 11.7, 1.1, size=15, color=MUTED, line_spacing=1.2)
feats = [
    ("Upload", "Drop a CSV/XLSX — validated, parsed, profiled, quality-scanned."),
    ("Ask", "Natural-language questions answered with KPIs, charts, and tables."),
    ("Verify", "Every number carries a receipt and a re-verifiable certificate."),
]
for i, (h, b) in enumerate(feats):
    x = 0.9 + i * 3.95
    rrect(s, x, 3.15, 3.65, 1.7, fill=LIGHT, line=BORDER, radius=0.07)
    text(s, h, x + 0.25, 3.32, 3.2, 0.4, size=17, color=VIOLET, bold=True, font=BODY)
    text(s, b, x + 0.25, 3.78, 3.2, 0.95, size=13, color=MUTED, line_spacing=1.1)
rrect(s, 0.9, 5.2, 11.55, 1.25, fill=DEEP, radius=0.06)
text(s, "Deterministic-first guarantee", 1.2, 5.38, 5, 0.4, size=14, color=EMERALD, bold=True)
text(s, "Every metric you see is computed in Pandas / scikit-learn. The LLM only narrates — it can never change a number. "
        "That is what makes DataVerse's output auditable and reproducible.",
     1.2, 5.74, 11.0, 0.7, size=13.5, color=ICE, line_spacing=1.15)
pagenum(s, 5)

# ============================================================ SLIDE 6 — TWO-AGENT ARCHITECTURE
s = slide(); bg(s, WHITE); kicker(s, "03  ·  Architecture"); title(s, "Exactly two agents — clear ownership")
text(s, "No sprawling agent swarm. A clean division of labour keeps results fast, reliable, and reproducible.",
     0.9, 1.78, 11.5, 0.5, size=14, color=MUTED)
# Agent A
rrect(s, 0.9, 2.45, 5.55, 3.5, fill=WHITE, line=VIOLET, line_w=1.5, radius=0.05)
circle(s, 1.2, 2.7, 0.6, VIOLET, "1", WHITE, 18)
text(s, "DatasetAgent", 1.95, 2.78, 4.3, 0.45, size=18, color=INK, bold=True, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
text(s, "DATA CORE", 1.95, 3.18, 4, 0.3, size=10.5, color=VIOLET, bold=True)
bullets(s, ["Secure CSV / XLSX upload & parsing", "Header normalisation & type profiling",
            "Quality scan: missing values, duplicates", "Persists the session (local or Supabase)"],
        1.2, 3.65, 5.1, 2.2, size=13.5)
# Agent B
rrect(s, 6.85, 2.45, 5.6, 3.5, fill=WHITE, line=INDIGO, line_w=1.5, radius=0.05)
circle(s, 7.15, 2.7, 0.6, INDIGO, "2", WHITE, 18)
text(s, "AnalystAgent", 7.9, 2.78, 4.3, 0.45, size=18, color=INK, bold=True, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
text(s, "ANALYTICS  ·  XAI  ·  TRUST", 7.9, 3.18, 4.4, 0.3, size=10.5, color=INDIGO, bold=True)
bullets(s, ["Semantic KPI mapping, EDA, trends", "Optional Ridge / RandomForest prediction",
            "SHAP-grounded explainable AI", "Provenance, certificate & narration"],
        7.15, 3.65, 5.1, 2.2, size=13.5)
text(s, "DatasetAgent  →  AnalystAgent  →  AnalysisPipeline  →  Report Composer/Generator",
     0.9, 6.15, 11.6, 0.5, size=13.5, color=INK, bold=True, align=PP_ALIGN.CENTER, font=BODY)
pagenum(s, 6)

# ============================================================ SLIDE 7 — SYSTEM ARCHITECTURE (layers)
s = slide(); bg(s, WHITE); kicker(s, "03  ·  System Architecture"); title(s, "A layered, deterministic pipeline")
layers = [
    ("Presentation", "Next.js 15 + React 19 chat UI — upload, KPIs, charts, verification panels", VIOLET),
    ("Application (FastAPI)", "Session-based REST API · two agents · async pipeline orchestration", INDIGO),
    ("Analytics (deterministic)", "semantic map → query plan → EDA / metrics / trends → Ridge/RF + SHAP", EMERALD),
    ("Trust layer", "provenance receipts · audit trail · reproducibility certificate · what-if", RGBColor(0x0E,0x76,0x90)),
    ("Persistence", "Local filesystem session store (optional Supabase / PostgreSQL)", SLATE),
]
y = 1.95
for h, b, ac in layers:
    rrect(s, 0.9, y, 11.55, 0.92, fill=WHITE, line=BORDER, radius=0.06)
    circle(s, 1.1, y + 0.21, 0.5, ac, "", WHITE)
    text(s, h, 1.8, y + 0.06, 3.7, 0.8, size=15, color=INK, bold=True, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
    text(s, b, 5.4, y + 0.06, 6.9, 0.8, size=12.5, color=MUTED, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
    y += 1.02
pagenum(s, 7)

# ============================================================ SLIDE 8 — THE VERIFIABILITY STACK
s = slide(); bg(s, DEEP); kicker(s, "04  ·  The Differentiator", color=EMERALD)
title(s, "The Verifiability Stack — “prove every number”", color=WHITE)
text(s, "Six layers of trust that competitors don't offer. Two are genuinely unique.", 0.7, 1.55, 11.5, 0.4, size=14, color=ICE)
feat = [
    ("Show-the-Math receipts", "Every KPI expands to its exact formula, source columns, row count & sample rows.", False),
    ("Full Audit Trail", "Receipts for KPIs, EDA, correlations, trends & model — downloadable as JSON.", False),
    ("Reproducibility Certificate", "SHA-256 of data + results; one-click re-verify proves numbers reproduce exactly.", True),
    ("Verified What-If Simulator", "Adjust a lever; KPIs recompute deterministically, each with its own receipt.", True),
    ("Data Quality Doctor", "Detects issues & applies one-click deterministic fixes with before/after.", False),
    ("Verified Report Export", "Self-contained HTML where every KPI embeds its receipt + the certificate.", False),
]
for i, (h, b, uq) in enumerate(feat):
    col, row = i % 3, i // 3
    x = 0.7 + col * 4.05; y = 2.15 + row * 2.15
    rrect(s, x, y, 3.8, 1.95, fill=RGBColor(0x2A,0x26,0x5F), radius=0.06)
    text(s, h, x + 0.25, y + 0.18, 3.3, 0.6, size=14.5, color=WHITE, bold=True, font=BODY)
    text(s, b, x + 0.25, y + 0.74, 3.35, 1.1, size=11.5, color=ICE, line_spacing=1.08)
    if uq:
        badge(s, "★ UNIQUE", x + 2.75, y + 0.16, EMERALD, 1.0)
pagenum(s, 8)

# ============================================================ SLIDE 9 — RECEIPTS + AUDIT
s = slide(); bg(s, WHITE); kicker(s, "04  ·  Trust in action"); title(s, "Show the Math — receipts on every number")
bullets(s, [
    ("Click any KPI.", "It expands into a receipt: the operation, plain-English formula, source columns, row count, and a sample of the contributing rows."),
    ("Computed, never invented.", "The receipt value is asserted equal to the metric value in tests — the guarantee is enforced by CI."),
    ("Full audit trail.", "The same receipts cover EDA stats, correlations, trends and model metrics — and export as one JSON audit log."),
], 0.9, 1.95, 6.2, 3.8, size=15)
# mock receipt card
rrect(s, 7.5, 2.0, 5.1, 4.2, fill=LIGHT, line=BORDER, radius=0.05)
text(s, "Total Sales", 7.8, 2.2, 4.5, 0.3, size=11, color=SLATE, bold=True)
text(s, "$2,248,662.62", 7.8, 2.5, 4.5, 0.55, size=24, color=INK, bold=True, font=HEAD)
badge(s, "✓ Show the math", 7.8, 3.15, VIOLET, 1.7)
rrect(s, 7.8, 3.65, 4.5, 2.35, fill=CARDBG, line=RGBColor(0xDD,0xD6,0xFE), radius=0.05)
text(s, "SUM of `sales` across 35,000 rows = 2,248,662.62", 8.0, 3.8, 4.1, 0.6, size=12, color=INK, bold=True)
text(s, "SUM    ·    sales", 8.0, 4.35, 4.1, 0.3, size=11, color=VIOLET, font="Consolas")
text(s, "sample rows:  1,200.50  ·  980.00  ·  875.25  ·  …", 8.0, 4.7, 4.1, 0.3, size=11, color=MUTED)
text(s, "✓  Verified deterministically from your rows", 8.0, 5.5, 4.1, 0.4, size=11.5, color=EMERALD, bold=True)
pagenum(s, 9)

# ============================================================ SLIDE 10 — REPRODUCIBILITY CERTIFICATE
s = slide(); bg(s, WHITE); kicker(s, "05  ·  Unique Feature #1", color=EMERALD)
title(s, "Reproducibility Certificate")
badge(s, "★ UNIQUE vs ALL existing tools", 9.2, 0.95, EMERALD, 3.3)
text(s, "Every analysis is issued a tamper-evident certificate: a SHA-256 fingerprint of your raw data and a SHA-256 "
        "fingerprint of the deterministic results. One click on “Re-verify” re-runs the pipeline on the raw rows and "
        "proves every number reproduces exactly — and detects if the data or the numbers were altered.",
     0.9, 1.95, 11.6, 1.2, size=15, color=MUTED, line_spacing=1.2)
steps = [("1", "Analyse", "Numbers + certificate (two fingerprints) generated"),
         ("2", "Share", "Certificate travels inside the exported report"),
         ("3", "Re-verify", "Anyone re-runs on the data → exact match or tamper flag")]
for i, (n, h, b) in enumerate(steps):
    x = 0.9 + i * 3.95
    rrect(s, x, 3.35, 3.65, 1.6, fill=LIGHT, line=BORDER, radius=0.06)
    circle(s, x + 0.22, 3.55, 0.5, EMERALD, n, WHITE, 15)
    text(s, h, x + 0.85, 3.55, 2.6, 0.4, size=15, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, b, x + 0.25, 4.1, 3.2, 0.8, size=12, color=MUTED, line_spacing=1.05)
rrect(s, 0.9, 5.25, 11.55, 1.2, fill=DEEP, radius=0.06)
text(s, "Live result:  “Reproduced exactly — 20 / 20 numbers match.”", 1.2, 5.42, 11, 0.4, size=15, color=EMERALD, bold=True)
text(s, "No LLM-based analyst (ChatGPT ADA, Julius, Hex) can do this — their outputs aren't reproducible, so there is nothing to certify.",
     1.2, 5.8, 11, 0.5, size=13, color=ICE, line_spacing=1.15)
pagenum(s, 10)

# ============================================================ SLIDE 11 — WHAT-IF + DOCTOR
s = slide(); bg(s, WHITE); kicker(s, "05  ·  Unique Feature #2  +  more", color=EMERALD)
title(s, "Verified What-If  &  Data Quality Doctor")
# What-if
rrect(s, 0.9, 1.95, 5.7, 4.4, fill=WHITE, line=BORDER, line_w=1.2, radius=0.05)
badge(s, "★ UNIQUE", 5.35, 2.1, EMERALD, 1.1)
text(s, "Verified What-If Simulator", 1.15, 2.1, 4.2, 0.4, size=16, color=INK, bold=True)
text(s, "Adjust a numeric lever; KPIs recompute deterministically — each hypothetical number keeps its own receipt.",
     1.15, 2.5, 5.2, 0.85, size=12.5, color=MUTED, line_spacing=1.1)
rows = [("KPI", "Baseline", "+10% price", True),
        ("Total Sales", "$2.25M", "$2.47M", False),
        ("Total Profit", "$336.9K", "$561.8K", False),
        ("Gross Margin", "14.98%", "22.7%", False)]
yy = 3.45
for kpi, b, sc, hd in rows:
    c = INK if hd else MUTED
    text(s, kpi, 1.2, yy, 2.2, 0.3, size=12.5, color=(SLATE if hd else INK), bold=hd)
    text(s, b, 3.4, yy, 1.5, 0.3, size=12.5, color=(SLATE if hd else MUTED), bold=hd)
    text(s, sc, 4.95, yy, 1.5, 0.3, size=12.5, color=(SLATE if hd else EMERALD), bold=True)
    yy += 0.45
text(s, "Deterministic & receipt-backed — reproducible what-if, not an LLM guess.", 1.15, 5.55, 5.3, 0.7, size=11.5, color=EMERALD, italic=True, line_spacing=1.1)
# Doctor
rrect(s, 6.75, 1.95, 5.7, 4.4, fill=WHITE, line=BORDER, line_w=1.2, radius=0.05)
text(s, "Data Quality Doctor", 7.0, 2.1, 4.5, 0.4, size=16, color=INK, bold=True)
text(s, "Auto-detects issues and applies one-click deterministic fixes with a before/after diff.",
     7.0, 2.5, 5.2, 0.7, size=12.5, color=MUTED, line_spacing=1.1)
bullets(s, [("Detects:", "duplicates, missing values, constant columns, numeric-as-text"),
            ("Fixes:", "impute (median / mode), drop column, drop duplicates, cast type"),
            ("Proves:", "re-analyses the cleaned data and shows what changed")],
        7.0, 3.25, 5.3, 2.0, size=12.5)
text(s, "Beats tools that only describe problems — DataVerse fixes them, deterministically.", 7.0, 5.55, 5.3, 0.7, size=11.5, color=VIOLET, italic=True, line_spacing=1.1)
pagenum(s, 11)

# ============================================================ SLIDE 12 — TECH STACK
s = slide(); bg(s, WHITE); kicker(s, "06  ·  Technology Stack"); title(s, "Built on a lean, deterministic stack")
rows = [
    ("Frontend", "Next.js 15, React 19, TypeScript, Tailwind CSS, motion"),
    ("Backend", "FastAPI, Uvicorn (ASGI), session-based REST API"),
    ("Data processing", "Pandas, NumPy, openpyxl"),
    ("ML / prediction", "scikit-learn (Ridge, RandomForest), SciPy, joblib"),
    ("Explainable AI", "SHAP (+ feature-importance fallback)"),
    ("Charts & reports", "Hand-rolled SVG charts, ReportLab (PDF/HTML)"),
    ("Optional LLM", "OpenAI / Google — narration polish only (offline Mock mode)"),
    ("Trust & storage", "SHA-256 certificate; local store (optional Supabase)"),
]
y = 1.95
for i, (k, v) in enumerate(rows):
    fill = LIGHT if i % 2 == 0 else WHITE
    rrect(s, 0.9, y, 11.55, 0.56, fill=fill, line=BORDER, line_w=0.75, radius=0.04)
    text(s, k, 1.15, y + 0.04, 3.3, 0.5, size=13.5, color=VIOLET, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, v, 4.5, y + 0.04, 7.8, 0.5, size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.6
text(s, "No LangChain · No PyCaret · No LIME — the live pipeline is deliberately deterministic and dependency-light.",
     0.9, 6.85, 11.6, 0.4, size=12.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
pagenum(s, 12)

# ============================================================ SLIDE 13 — DATASET & TESTING
s = slide(); bg(s, WHITE); kicker(s, "07  ·  Dataset & Testing"); title(s, "Evaluated on a 35,000-row retail dataset")
rrect(s, 0.9, 1.95, 5.6, 4.35, fill=DEEP, radius=0.05)
text(s, "retail_mart_processed_v1.csv", 1.2, 2.15, 5.1, 0.4, size=15, color=WHITE, bold=True)
for i, (a, b) in enumerate([("35,000", "rows"), ("21", "columns"), ("8", "KPI dimensions"), ("1.0", "data-quality score")]):
    x = 1.2 + (i % 2) * 2.6; yy = 2.7 + (i // 2) * 1.0
    text(s, a, x, yy, 2.4, 0.5, size=27, color=EMERALD, bold=True, font=HEAD)
    text(s, b, x, yy + 0.5, 2.4, 0.3, size=11.5, color=ICE)
text(s, "Region · Category · Customer Type · Payment · Online/Offline · time fields",
     1.2, 5.55, 5.1, 0.7, size=12, color=ICE, line_spacing=1.15)
text(s, "Tests performed", 6.95, 1.95, 5, 0.4, size=15, color=INK, bold=True)
bullets(s, ["CSV upload via REST API", "Dataset profiling & EDA pipeline",
            "KPI generation (sales, profit, margin)", "Region & category breakdowns",
            "ML prediction + SHAP explainability", "Data Quality Doctor + cleaning",
            "Reproducibility certificate + re-verify", "HTML / PDF verified report export"],
        6.95, 2.4, 5.5, 3.9, size=13.5, gap=5)
pagenum(s, 13)

# ============================================================ SLIDE 14 — RESULTS (KPIs)
s = slide(); bg(s, WHITE); kicker(s, "08  ·  Results"); title(s, "Business insights — every figure verifiable")
stats = [("$2,248,662.62", "Total Sales"), ("$336,938.93", "Total Profit"),
         ("14.98%", "Profit Margin"), ("69,546", "Units Sold")]
for i, (a, b) in enumerate(stats):
    x = 0.9 + i * 2.95
    rrect(s, x, 1.95, 2.7, 1.5, fill=LIGHT, line=BORDER, radius=0.06)
    text(s, a, x + 0.15, 2.12, 2.4, 0.6, size=21, color=VIOLET, bold=True, font=HEAD)
    text(s, b, x + 0.15, 2.78, 2.4, 0.4, size=12, color=SLATE)
bullets(s, [
    ("Best region:", "Region 2 — $1.28M in sales (prioritise stock & marketing)."),
    ("Best category:", "Category 1 — highest sales and profit."),
    ("Top store:", "Store 27 leads overall sales."),
    ("Channel:", "Online 15.03% vs Offline 14.94% margin — near-identical."),
    ("Discount vs profit:", "correlation -0.0739 — weak; discounts don't drive profit."),
], 0.9, 3.75, 11.6, 2.6, size=14.5, gap=7)
pagenum(s, 14)

# ============================================================ SLIDE 15 — VERIFICATION RESULTS
s = slide(); bg(s, WHITE); kicker(s, "08  ·  Engineering Verification"); title(s, "Proven, not just claimed")
cards3 = [("125", "backend tests passing", "pytest — including the trust test that asserts every receipt equals its metric", EMERALD),
          ("20 / 20", "numbers re-verified live", "certificate re-run reproduced every deterministic number exactly", VIOLET),
          ("0", "console errors", "frontend lint + production build green; verified in a real browser", INDIGO)]
for i, (a, b, c, ac) in enumerate(cards3):
    x = 0.9 + i * 3.95
    rrect(s, x, 1.95, 3.65, 2.5, fill=WHITE, line=BORDER, line_w=1.2, radius=0.05)
    text(s, a, x + 0.25, 2.15, 3.2, 0.7, size=34, color=ac, bold=True, font=HEAD)
    text(s, b, x + 0.25, 2.95, 3.2, 0.4, size=13.5, color=INK, bold=True)
    text(s, c, x + 0.25, 3.35, 3.2, 1.0, size=11.5, color=MUTED, line_spacing=1.1)
rrect(s, 0.9, 4.7, 11.55, 1.55, fill=LIGHT, line=BORDER, radius=0.06)
text(s, "Performance", 1.2, 4.85, 4, 0.35, size=13, color=VIOLET, bold=True)
bullets(s, [("Database insert:", "< 100 ms"), ("Full 35K-row analysis:", "< 500 ms"),
            ("Offline Mock mode:", "full pipeline runs with zero API keys")],
        1.2, 5.2, 11.0, 1.0, size=13.5, gap=3)
pagenum(s, 15)

# ============================================================ SLIDE 16 — COMPARISON
s = slide(); bg(s, WHITE); kicker(s, "09  ·  How We Compare"); title(s, "DataVerse vs existing tools")
headers = ["Capability", "ChatGPT ADA", "Julius AI", "Tableau", "DataVerse"]
data = [
    ["Natural-language Q&A", "✓", "✓", "—", "✓"],
    ["Deterministic numbers", "—", "—", "✓", "✓"],
    ["Per-number receipts", "—", "—", "—", "✓"],
    ["Reproducibility certificate", "—", "—", "—", "✓"],
    ["Verified what-if", "—", "—", "—", "✓"],
    ["One-click data cleaning", "partial", "partial", "—", "✓"],
]
colx = [0.9, 4.7, 6.7, 8.7, 10.7]
colw = [3.8, 2.0, 2.0, 2.0, 1.85]
y = 1.95
# header row
rrect(s, 0.9, y, 11.65, 0.55, fill=DEEP, radius=0.04)
for j, h in enumerate(headers):
    text(s, h, colx[j] + 0.1, y + 0.02, colw[j], 0.5, size=12.5, color=(EMERALD if j == 4 else WHITE), bold=True, anchor=MSO_ANCHOR.MIDDLE, align=(PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER))
y += 0.6
for i, row in enumerate(data):
    fill = LIGHT if i % 2 == 0 else WHITE
    rrect(s, 0.9, y, 11.65, 0.62, fill=fill, line=BORDER, line_w=0.75, radius=0.03)
    rrect(s, colx[4], y, colw[4] + 0.2, 0.62, fill=RGBColor(0xEC,0xFD,0xF5), radius=0.03)
    for j, val in enumerate(row):
        col = INK if j == 0 else (EMERALD if val == "✓" else (SLATE if val in ("—", "partial") else INK))
        text(s, val, colx[j] + 0.1, y + 0.04, colw[j], 0.55, size=12.5, color=col,
             bold=(j == 0 or j == 4), anchor=MSO_ANCHOR.MIDDLE, align=(PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER))
    y += 0.67
pagenum(s, 16)

# ============================================================ SLIDE 17 — CONCLUSION
s = slide(); bg(s, WHITE); kicker(s, "10  ·  Conclusion"); title(s, "What we achieved — and what's next")
text(s, "Achieved", 0.95, 1.95, 5, 0.4, size=16, color=VIOLET, bold=True)
bullets(s, ["A two-agent, deterministic BI analyst with NL chat",
            "A six-layer verifiability stack (2 features unique)",
            "Reproducibility certificate + verified what-if",
            "125 passing tests; live, browser-verified",
            "Polished, interactive Next.js workspace"],
        0.95, 2.4, 5.7, 3.6, size=14, gap=8)
text(s, "Future work", 7.0, 1.95, 5, 0.4, size=16, color=INDIGO, bold=True)
bullets(s, ["Public, link-based third-party re-verification",
            "Drift detection across dataset versions",
            "Provenance for deeper statistical tests",
            "Multi-file joins & larger datasets",
            "Role-based access & team workspaces"],
        7.0, 2.4, 5.6, 3.6, size=14, gap=8)
rrect(s, 0.9, 6.2, 11.55, 0.75, fill=DEEP, radius=0.07)
text(s, "DataVerse AI: the data analyst that doesn't just give you numbers — it proves them.",
     1.0, 6.32, 11.3, 0.5, size=14.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
pagenum(s, 17)

# ============================================================ SLIDE 18 — REFERENCES
s = slide(); bg(s, WHITE); kicker(s, "References"); title(s, "References")
refs = [
    "[1] V. V. Meduri et al., “BIREC: Guided Data Analysis for Conversational Business Intelligence,” arXiv:2105.00467, 2021.",
    "[2] J. Lian et al., “ChatBI: Towards Natural Language to Complex Business Intelligence SQL,” arXiv:2405.00527, 2024.",
    "[3] E. Cambria et al., “XAI meets LLMs: A Survey of the Relation between Explainable AI and LLMs,” arXiv:2407.15248, 2024.",
    "[4] “Survey on Explainable AI: Approaches, Limitations and Applications,” Human-Centric Intelligent Systems, 2023.",
    "[5] L. Shi et al., “A Survey on Employing LLMs for Text-to-SQL Tasks,” ACM Comput. Surv., 58(2), 2026, doi:10.1145/3737873.",
    "[6] S. M. Lundberg & S.-I. Lee, “A Unified Approach to Interpreting Model Predictions (SHAP),” NeurIPS, 2017.",
]
tb = s.shapes.add_textbox(Inches(0.9), Inches(1.95), Inches(11.6), Inches(4.6))
tf = tb.text_frame; tf.word_wrap = True
for i, r in enumerate(refs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(11); p.line_spacing = 1.1
    run = p.add_run(); run.text = r; run.font.size = Pt(13.5); run.font.color.rgb = MUTED; run.font.name = BODY
pagenum(s, 18)

# ============================================================ SLIDE 19 — THANK YOU
s = slide(); bg(s, DEEP)
circle(s, 6.07, 2.15, 1.2, VIOLET, "DV", WHITE, 34)
text(s, "Thank You", 0, 3.55, 13.333, 1.0, size=52, color=WHITE, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
text(s, "DataVerse AI — The Verifiable AI Data Analyst", 0, 4.6, 13.333, 0.5, size=20, color=ICE, align=PP_ALIGN.CENTER, font=HEAD, italic=True)
text(s, "Muhammad Muavia  ·  Hasnain Ali  ·  M. Shahan Subhani        Supervisor: Ms. Tayyaba Shehzad",
     0, 5.45, 13.333, 0.4, size=13.5, color=RGBColor(0xCB,0xD5,0xE1), align=PP_ALIGN.CENTER)
text(s, "Questions & Discussion", 0, 6.0, 13.333, 0.4, size=14, color=EMERALD, align=PP_ALIGN.CENTER, bold=True)

OUT = r"C:/Users/mouav/Downloads/FYP_Defense_DataVerse_Final.pptx"
prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides))
