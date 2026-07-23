from __future__ import annotations

import math
import re
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.text.paragraph import Paragraph
from docx.shared import Inches, Pt
from pptx.dml.color import RGBColor as PptRGBColor
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
SOURCE_DOCX = ROOT / "tmp" / "pdfs" / "report-source.docx"
SOURCE_PPTX = Path(r"C:\Users\mouav\OneDrive\Desktop\FYP P2 Defense (Final v2).pptx")
REPORT_OUT = ROOT / "output" / "pdf" / "DataVerse_AI_Final_Report_Revised_2026.docx"
PPT_OUT = ROOT / "output" / "slides" / "FYP_P2_Defense_Revised_2026.pptx"
ARCH_OUT = ROOT / "output" / "architecture"
UI_DIR = ROOT / "output" / "playwright"


FONT_REGULAR = Path(r"C:\Windows\Fonts\segoeui.ttf")
FONT_SEMIBOLD = Path(r"C:\Windows\Fonts\seguisb.ttf")
FONT_BOLD = Path(r"C:\Windows\Fonts\segoeuib.ttf")


def font(size: int, *, bold: bool = False, semibold: bool = False) -> ImageFont.FreeTypeFont:
    selected = FONT_BOLD if bold else FONT_SEMIBOLD if semibold else FONT_REGULAR
    return ImageFont.truetype(str(selected), size=size)


def wrap(draw: ImageDraw.ImageDraw, text: str, selected_font: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    words = text.split()
    lines: list[str] = []
    current = ""
    for word in words:
        trial = f"{current} {word}".strip()
        if draw.textbbox((0, 0), trial, font=selected_font)[2] <= max_width:
            current = trial
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def centered_text(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    title: str,
    subtitle: str = "",
    *,
    title_size: int = 34,
    subtitle_size: int = 25,
    color: str = "#0F172A",
) -> None:
    x1, y1, x2, y2 = box
    title_font = font(title_size, semibold=True)
    subtitle_font = font(subtitle_size)
    title_lines = wrap(draw, title, title_font, x2 - x1 - 38)
    subtitle_lines = wrap(draw, subtitle, subtitle_font, x2 - x1 - 42) if subtitle else []
    title_h = len(title_lines) * (title_size + 7)
    subtitle_h = len(subtitle_lines) * (subtitle_size + 6)
    gap = 10 if subtitle_lines else 0
    y = y1 + ((y2 - y1) - title_h - subtitle_h - gap) // 2
    for line in title_lines:
        bbox = draw.textbbox((0, 0), line, font=title_font)
        draw.text(((x1 + x2 - (bbox[2] - bbox[0])) / 2, y), line, font=title_font, fill=color)
        y += title_size + 7
    y += gap
    for line in subtitle_lines:
        bbox = draw.textbbox((0, 0), line, font=subtitle_font)
        draw.text(((x1 + x2 - (bbox[2] - bbox[0])) / 2, y), line, font=subtitle_font, fill="#475569")
        y += subtitle_size + 6


def box(
    draw: ImageDraw.ImageDraw,
    coords: tuple[int, int, int, int],
    title: str,
    subtitle: str,
    *,
    fill: str,
    outline: str,
) -> None:
    draw.rounded_rectangle(coords, radius=24, fill=fill, outline=outline, width=3)
    centered_text(draw, coords, title, subtitle)


def arrow_head(draw: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int], color: str, width: int = 5) -> None:
    angle = math.atan2(end[1] - start[1], end[0] - start[0])
    length = 18
    spread = math.pi / 6
    p1 = (end[0] - length * math.cos(angle - spread), end[1] - length * math.sin(angle - spread))
    p2 = (end[0] - length * math.cos(angle + spread), end[1] - length * math.sin(angle + spread))
    draw.polygon([end, p1, p2], fill=color)


def connector(
    draw: ImageDraw.ImageDraw,
    points: list[tuple[int, int]],
    *,
    color: str = "#475569",
    label: str = "",
    bidirectional: bool = False,
    dashed: bool = False,
) -> None:
    if dashed:
        for a, b in zip(points, points[1:]):
            length = max(abs(b[0] - a[0]), abs(b[1] - a[1]))
            steps = max(1, length // 24)
            for i in range(0, steps, 2):
                t1, t2 = i / steps, min(1, (i + 1) / steps)
                s = (int(a[0] + (b[0] - a[0]) * t1), int(a[1] + (b[1] - a[1]) * t1))
                e = (int(a[0] + (b[0] - a[0]) * t2), int(a[1] + (b[1] - a[1]) * t2))
                draw.line([s, e], fill=color, width=5)
    else:
        draw.line(points, fill=color, width=5, joint="curve")
    arrow_head(draw, points[-2], points[-1], color)
    if bidirectional:
        arrow_head(draw, points[1], points[0], color)
    if label:
        longest_start, longest_end = max(
            zip(points, points[1:]),
            key=lambda segment: math.dist(segment[0], segment[1]),
        )
        mid = ((longest_start[0] + longest_end[0]) // 2, (longest_start[1] + longest_end[1]) // 2)
        label_font = font(22, semibold=True)
        bbox = draw.textbbox((0, 0), label, font=label_font)
        pad = 8
        label_box = (mid[0] - (bbox[2] - bbox[0]) // 2 - pad, mid[1] - 36, mid[0] + (bbox[2] - bbox[0]) // 2 + pad, mid[1] - 4)
        draw.rounded_rectangle(label_box, radius=7, fill="#FFFFFF")
        draw.text((label_box[0] + pad, label_box[1] + 2), label, font=label_font, fill="#334155")


def make_architecture_diagram(path: Path) -> None:
    img = Image.new("RGB", (2400, 1350), "#F8FAFC")
    draw = ImageDraw.Draw(img)
    draw.text((80, 55), "DataVerse AI System Architecture", font=font(58, bold=True), fill="#0F172A")
    draw.text((82, 130), "Simple view of the components used by the application", font=font(30), fill="#475569")

    user = (60, 395, 360, 635)
    frontend = (420, 395, 760, 635)
    backend = (820, 395, 1160, 635)
    agents = (1220, 395, 1580, 635)
    analytics = (1640, 395, 1980, 635)
    output = (2040, 395, 2340, 635)
    supabase = (760, 830, 1400, 1110)

    box(draw, user, "User", "Signs in, uploads data, asks questions", fill="#FFFFFF", outline="#8B5CF6")
    box(draw, frontend, "Next.js frontend", "Interface, session, charts, reports", fill="#FFFFFF", outline="#8B5CF6")
    box(draw, backend, "FastAPI backend", "Secure API, validation, routing", fill="#FFFFFF", outline="#2563EB")
    box(draw, agents, "Two analysis agents", "DatasetAgent and AnalystAgent", fill="#FFFFFF", outline="#2563EB")
    box(draw, analytics, "Analytics engine", "Pandas, scikit-learn, SHAP, reports", fill="#FFFFFF", outline="#2563EB")
    box(draw, output, "User output", "Verified insights, charts, and downloadable reports", fill="#FFF7ED", outline="#EA580C")
    box(draw, supabase, "Supabase", "Auth, PostgreSQL, and Storage for accounts, metadata, datasets, and reports", fill="#ECFDF5", outline="#059669")

    connector(draw, [(360, 515), (420, 515)])
    connector(draw, [(760, 515), (820, 515)])
    connector(draw, [(1160, 515), (1220, 515)])
    connector(draw, [(1580, 515), (1640, 515)])
    connector(draw, [(1980, 515), (2040, 515)])
    connector(draw, [(990, 635), (990, 830)], bidirectional=True, label="auth, metadata, files")
    draw.text((720, 345), "HTTPS + JWT", font=font(22, semibold=True), fill="#334155")

    draw.text((600, 1230), "Requests move left to right; verified results are presented in the user output box.", font=font(28, semibold=True), fill="#334155")
    img.save(path, quality=95)


def make_flow_diagram(path: Path) -> None:
    img = Image.new("RGB", (2400, 1350), "#FFFFFF")
    draw = ImageDraw.Draw(img)
    draw.text((80, 45), "DataVerse AI Request and Data Flow", font=font(54, bold=True), fill="#0F172A")
    draw.text((82, 112), "The authentication branch and the analysis branch use separate Supabase services", font=font(28), fill="#475569")

    main_x1, main_x2 = 760, 1640
    steps = [
        ((main_x1, 190, main_x2, 315), "1. Browser and Next.js UI", "User signs in, uploads a CSV / Excel file, or asks a question"),
        ((main_x1, 380, main_x2, 505), "2. FastAPI API boundary", "Validates input, verifies the JWT, checks ownership, applies rate limits"),
        ((main_x1, 570, main_x2, 695), "3. SessionService / orchestrator", "Loads session context and selects the upload or analysis workflow"),
        ((main_x1, 760, main_x2, 885), "4. DatasetAgent then AnalystAgent", "Profiles the file, plans the request, and invokes deterministic modules"),
        ((main_x1, 950, main_x2, 1075), "5. Pandas, scikit-learn, and XAI", "Computes metrics, EDA, models, explanations, and report content"),
        ((main_x1, 1140, main_x2, 1265), "6. Verified response to the UI", "Returns JSON, charts, receipts, and report links to the signed-in user"),
    ]
    for coords, title, subtitle in steps:
        box(draw, coords, title, subtitle, fill="#F8FAFC", outline="#3B82F6")
    for first, second in zip(steps, steps[1:]):
        connector(draw, [((first[0][0] + first[0][2]) // 2, first[0][3]), ((second[0][0] + second[0][2]) // 2, second[0][1])])

    auth = (1775, 340, 2290, 535)
    store = (1775, 710, 2290, 955)
    llm = (110, 760, 625, 955)
    box(draw, auth, "Supabase Auth", "Email confirmation and password handling. Returns signed JWTs to the client through FastAPI.", fill="#ECFDF5", outline="#10B981")
    box(draw, store, "Supabase PostgreSQL + Storage", "Metadata in PostgreSQL. Dataset and report files in private storage buckets.", fill="#ECFDF5", outline="#10B981")
    box(draw, llm, "Optional LLM provider", "May refine intent and wording. Numeric results remain deterministic.", fill="#FFF7ED", outline="#F59E0B")

    connector(draw, [(1640, 442), (1775, 442)], bidirectional=True, label="sign-up / login / refresh")
    connector(draw, [(1640, 635), (1710, 635), (1710, 805), (1775, 805)], bidirectional=True, label="metadata")
    connector(draw, [(1640, 1012), (1715, 1012), (1715, 900), (1775, 900)], label="files and results")
    connector(draw, [(760, 820), (625, 820)], bidirectional=True, dashed=True, label="intent / narration")

    security_note = (95, 245, 620, 610)
    draw.rounded_rectangle(security_note, radius=24, fill="#F8FAFC", outline="#64748B", width=3)
    draw.text((125, 275), "Security boundary", font=font(29, semibold=True), fill="#0F172A")
    security_lines = [
        "The public anon key is used only for",
        "Supabase Auth REST calls.",
        "",
        "Private database and storage work",
        "stays on the backend service role.",
        "",
        "Supabase Auth hashes passwords;",
        "DataVerse AI never stores them.",
    ]
    y_note = 325
    for line in security_lines:
        draw.text((125, y_note), line, font=font(22), fill="#475569")
        y_note += 29
    img.save(path, quality=95)


PHRASE_REPLACEMENTS = [
    ("The purpose of this chapter is to", "This chapter"),
    ("This chapter presents", "This chapter records"),
    ("This chapter provides", "This chapter gives"),
    ("The proposed system", "The implemented application"),
    ("the proposed system", "the implemented application"),
    ("DataVerse AI is designed to", "We built DataVerse AI to"),
    ("This study aims to", "This project examines how to"),
    ("This research aims to", "This project examines how to"),
    ("in order to", "to"),
    ("utilizes", "uses"),
    ("utilize", "use"),
    ("leverages", "uses"),
    ("leverage", "use"),
    ("facilitates", "supports"),
    ("allows users to", "lets users"),
    ("enables users to", "lets users"),
    ("Moreover,", "Also,"),
    ("Furthermore,", "In addition,"),
    ("Additionally,", "The project also"),
    ("It is important to note that ", ""),
    ("plays a crucial role in", "is used for"),
    ("plays a vital role in", "is used for"),
    ("a comprehensive", "a detailed"),
    ("seamlessly", ""),
    ("robust", "reliable"),
    ("user-friendly", "clear"),
    ("as well as", "and"),
    ("demonstrates", "shows"),
    ("The results demonstrate", "The results show"),
    ("This approach ensures that", "This keeps"),
    ("The system", "The application"),
    ("the system", "the application"),
    ("the platform", "the product"),
    ("The platform", "The product"),
    ("provides users with", "gives users"),
    ("provides", "offers"),
    ("ensures", "makes sure"),
    ("ensure", "make sure"),
    ("enhances", "improves"),
    ("enhance", "improve"),
    ("valuable insights", "useful findings"),
    ("actionable insights", "findings that can guide a decision"),
    ("insights", "findings"),
    ("various", "several"),
    ("significant", "notable"),
    ("effectively", ""),
    ("efficiently", ""),
    ("Overall,", "Taken together,"),
    ("In conclusion,", "In summary,"),
    ("Due to the fact that", "Because"),
    ("a wide range of", "several"),
    ("has the ability to", "can"),
    ("is capable of", "can"),
    ("It can be observed that", "The results show that"),
    ("It was observed that", "We observed that"),
    ("This chapter documents", "Here we document"),
    ("This chapter reports", "Here we report"),
    ("This chapter covered", "This chapter recorded"),
    ("This chapter framed", "This chapter set out"),
    ("This chapter outlines", "This chapter lists"),
    ("This chapter explains", "This chapter walks through"),
]


def clean_human_text(text: str) -> str:
    original = text
    text = text.replace("—", ",").replace("–", "-")
    for old, new in PHRASE_REPLACEMENTS:
        text = text.replace(old, new)
    text = re.sub(r"\s+", " ", text).strip()
    text = text.replace(" ,", ",").replace("..", ".")
    if ";" in text:
        text = text.replace("; ", ". ")
    if text.startswith("The application "):
        text = "In the implemented version, the application " + text[len("The application "):]
    if text.startswith("The architecture "):
        text = "In the implemented architecture, " + text[len("The architecture "):]
    return text if text else original


def replace_paragraph_text(paragraph: Paragraph, text: str) -> None:
    if paragraph.text == text:
        return
    for run in paragraph.runs:
        run._element.getparent().remove(run._element)
    paragraph.add_run(text)


def find_paragraph(doc: Document, exact_text: str) -> Paragraph:
    for paragraph in doc.paragraphs:
        if paragraph.text.strip() == exact_text:
            return paragraph
    raise KeyError(exact_text)


def paragraph_before(doc: Document, paragraph: Paragraph) -> Paragraph:
    paragraphs = doc.paragraphs
    index = next(i for i, candidate in enumerate(paragraphs) if candidate._p is paragraph._p)
    return paragraphs[index - 1]


def paragraph_after(doc: Document, paragraph: Paragraph) -> Paragraph:
    paragraphs = doc.paragraphs
    index = next(i for i, candidate in enumerate(paragraphs) if candidate._p is paragraph._p)
    return paragraphs[index + 1]


def set_body_format(paragraph: Paragraph) -> None:
    paragraph.style = "Normal"
    paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    for run in paragraph.runs:
        run.bold = False
        run.italic = False
        run.font.size = Pt(11)


def set_picture(paragraph: Paragraph, image_path: Path, width_inches: float) -> None:
    for child in list(paragraph._p):
        if child.tag.endswith("}r") or child.tag.endswith("}hyperlink"):
            paragraph._p.remove(child)
    paragraph.add_run().add_picture(str(image_path), width=Inches(width_inches))
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER


def delete_paragraph(paragraph: Paragraph) -> None:
    parent = paragraph._element.getparent()
    parent.remove(paragraph._element)


def insert_caption_before_table(table, text: str) -> None:
    xml_p = OxmlElement("w:p")
    table._tbl.addprevious(xml_p)
    paragraph = Paragraph(xml_p, table._parent)
    paragraph.style = "Normal"
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run(text)
    run.bold = True
    run.font.size = Pt(10)


def update_docx() -> tuple[int, int]:
    doc = Document(SOURCE_DOCX)

    changed = 0
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text or text.startswith("[") or text.lower().startswith(("figure ", "table ")):
            continue
        if paragraph.style.name == "Normal" and len(text) >= 140:
            revised = clean_human_text(text)
            if revised != text:
                replace_paragraph_text(paragraph, revised)
                changed += 1
        elif "—" in text or "–" in text:
            replace_paragraph_text(paragraph, clean_human_text(text))
            changed += 1

    targeted = {
        "The architecture describes a conversational AI application.": (
            "Figure 1 shows the main components used by DataVerse AI. The user works through the Next.js frontend, which exchanges HTTPS requests and JWT-authenticated responses with FastAPI. FastAPI routes uploads and questions to DatasetAgent and AnalystAgent. These agents call the deterministic Pandas, scikit-learn, SHAP, and reporting modules. Supabase provides authentication, PostgreSQL metadata storage, and private file storage for datasets and reports."
        ),
        "The flowchart traces a request across the whole platform": (
            "Figure 20 separates the account flow from the analysis flow. Sign-up, login, refresh, and email confirmation go to Supabase Auth through FastAPI. Dataset uploads and questions go to SessionService, then to the two agents and the deterministic analytics modules. Metadata is written to Supabase PostgreSQL, while dataset and report files are written to private Storage buckets. The response travels back through FastAPI to the Next.js interface."
        ),
        "The interface is built around simplicity": (
            "The current interface uses four public screens: the landing page, the feature overview, secure login, and secure sign-up. Guest access has been removed. New accounts require a 12-character password with mixed character classes and must be confirmed through the email link sent by Supabase Auth."
        ),
    }
    for paragraph in doc.paragraphs:
        raw = paragraph.text.strip()
        for prefix, replacement in targeted.items():
            if raw.startswith(prefix):
                replace_paragraph_text(paragraph, replacement)
                changed += 1

    for paragraph in doc.paragraphs:
        if "register, sign in or continue as guest" in paragraph.text:
            replace_paragraph_text(paragraph, paragraph.text.replace("register, sign in or continue as guest", "register, confirm an email address, and sign in"))
            changed += 1
        if "guest entry" in paragraph.text.lower():
            replace_paragraph_text(paragraph, re.sub("guest entry", "confirmed account access", paragraph.text, flags=re.I))
            changed += 1
        if "207" in paragraph.text:
            replace_paragraph_text(paragraph, paragraph.text.replace("207", "214"))
            changed += 1

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    if "207" in paragraph.text:
                        replace_paragraph_text(paragraph, paragraph.text.replace("207", "214"))

    architecture = ARCH_OUT / "dataverse-production-architecture.png"
    flow = ARCH_OUT / "dataverse-request-data-flow.png"
    for caption_text, image_path in [
        ("Figure 1: System Architecture Diagram", architecture),
        ("Figure 19: Architecture Diagram", architecture),
        ("Figure 20: System Flowchart", flow),
    ]:
        caption = find_paragraph(doc, caption_text)
        set_picture(paragraph_before(doc, caption), image_path, 6.6)

    figure1_caption = find_paragraph(doc, "Figure 1: System Architecture Diagram")
    figure1_body = paragraph_after(doc, figure1_caption)
    replace_paragraph_text(
        figure1_body,
        "Figure 1 shows the main components used by DataVerse AI. The user works through the Next.js frontend, which exchanges HTTPS requests and JWT-authenticated responses with FastAPI. FastAPI routes uploads and questions to DatasetAgent and AnalystAgent. These agents call the deterministic Pandas, scikit-learn, SHAP, and reporting modules. Supabase provides authentication, PostgreSQL metadata storage, and private file storage for datasets and reports.",
    )
    set_body_format(figure1_body)

    for paragraph in doc.paragraphs:
        if paragraph.text.startswith("Figure 20 separates the account flow"):
            set_body_format(paragraph)
            break

    component_updates = {
        "User Web / Mobile Client": (
            "Browser and Next.js UI",
            "The signed-in user uploads a CSV or Excel file and asks questions in the Next.js interface. The browser keeps only the session state and sends HTTPS requests through the same-origin backend proxy.",
        ),
        "FastAPI REST API": (
            "FastAPI API Boundary",
            "FastAPI validates request data, verifies the Supabase JWT, checks resource ownership, applies rate limits, and returns security headers before any analysis or storage operation runs.",
        ),
        "API Routes (Intent / Query Router)": (
            "SessionService and Orchestrator",
            "SessionService loads the active session, messages, and dataset metadata. It then selects the upload, question, analysis, or report workflow and records the resulting agent run.",
        ),
        "Orchestrator (Workflow Manager)": (
            "DatasetAgent",
            "DatasetAgent checks the file extension and size, parses the upload, normalizes column names, profiles the DataFrame, and records data-quality results before analysis begins.",
        ),
        "Agent Layer (DatasetAgent + AnalystAgent)": (
            "AnalystAgent",
            "AnalystAgent turns the request into an analysis plan and calls the deterministic modules for KPIs, EDA, prediction, explanation, and report composition. It does not invent numeric results.",
        ),
        "Local Deep Analysis": (
            "Deterministic Analytics Engine",
            "Pandas computes the statistics, business metrics, trends, correlations, and outliers. Scikit-learn trains supported models, while SHAP or feature importance explains model behavior.",
        ),
        "Pandas DataFrame": (
            "Supabase Auth",
            "Supabase Auth owns password hashing, email confirmation, login, token refresh, and JWT issuance. DataVerse AI never stores raw passwords in its database or application code.",
        ),
        "Scikit-learn": (
            "Supabase PostgreSQL",
            "The backend service role writes session, message, dataset, agent-run, and report metadata to PostgreSQL. Row ownership is always checked against the authenticated user.",
        ),
        "SHAP / Feature Importance": (
            "Supabase Storage",
            "Private Storage buckets hold the uploaded dataset files and generated reports. Only the backend uses the service-role credential for these operations.",
        ),
        "Supabase (PostgreSQL + Storage)": (
            "Optional LLM Provider",
            "An optional language model may refine intent or improve wording. It is outside the deterministic calculation path and therefore cannot supply a KPI, model metric, or chart value.",
        ),
        "Final Outputs (Tables, Datasets, Queries, Agent Runs, Reports)": (
            "Verified Response and Reports",
            "FastAPI returns JSON results, charts, receipts, and report links to the Next.js interface. The response is shown only in the authenticated user's workspace.",
        ),
    }
    for old_heading, (new_heading, description) in component_updates.items():
        heading = find_paragraph(doc, old_heading)
        body = paragraph_after(doc, heading)
        replace_paragraph_text(heading, new_heading)
        replace_paragraph_text(body, description)
        set_body_format(body)

    screen_map = [
        ("Figure 21: Landing Page Screen", UI_DIR / "latest-landing.png", "Figure 21: Current Landing Page"),
        ("Figure 22: Login Screen", UI_DIR / "latest-features.png", "Figure 22: Current Features Page"),
        ("Figure 23: Sign-in Screen", UI_DIR / "latest-login.png", "Figure 23: Secure Login Page"),
        ("Figure 24: Sign-up Screen", UI_DIR / "latest-signup.png", "Figure 24: Secure Sign-up Page with Email Confirmation"),
    ]
    for old_caption, image_path, new_caption in screen_map:
        caption = find_paragraph(doc, old_caption)
        set_picture(paragraph_before(doc, caption), image_path, 6.55)
        replace_paragraph_text(caption, new_caption)

    for removable in ["Figure 25: Dashboard Screen", "Figure 26: Dashboard Overview"]:
        caption = find_paragraph(doc, removable)
        image_paragraph = paragraph_before(doc, caption)
        delete_paragraph(image_paragraph)
        delete_paragraph(caption)
    try:
        dashboard_label = find_paragraph(doc, "Dashboard")
        delete_paragraph(dashboard_label)
    except KeyError:
        pass
    find_paragraph(doc, "Landing Page:").text = "Current public pages"
    find_paragraph(doc, "Sign-in and Sign-up").text = "Secure account access"

    missing_captions = [
        "Table 2.1: Comparative Review of Related BI and Conversational Analytics Systems",
        "Table 3.1: Functional Requirements",
        "Table 3.2: Non-functional Requirements",
        "Table 3.3: Supabase Database Schema",
    ]
    for table, caption in zip(doc.tables[:4], missing_captions):
        insert_caption_before_table(table, caption)

    caption_pattern = re.compile(r"^(Figure|Table)\s+(?:\d+(?:\.\d+)?):", re.IGNORECASE)
    for paragraph in doc.paragraphs:
        if caption_pattern.match(paragraph.text.strip()):
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in paragraph.runs:
                run.bold = True
                run.font.size = Pt(10)

    props = doc.core_properties
    props.title = "DataVerse AI Final Report, Revised 2026"
    props.subject = "Updated architecture, Supabase flow, current UI, table captions, and revised prose"
    props.comments = "Revised against the implemented repository flow."
    REPORT_OUT.parent.mkdir(parents=True, exist_ok=True)
    doc.save(REPORT_OUT)
    return changed, len(doc.tables)


def iter_text_frames(shape) -> Iterable:
    if getattr(shape, "has_text_frame", False):
        yield shape.text_frame
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                yield cell.text_frame
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            yield from iter_text_frames(child)


def replace_ppt_picture(slide, old_shape, image_path: Path) -> None:
    left, top, width, height = old_shape.left, old_shape.top, old_shape.width, old_shape.height
    old_shape._element.getparent().remove(old_shape._element)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def update_pptx() -> int:
    prs = Presentation(SOURCE_PPTX)
    replacements = {
        "Supabase Authentication (Secure Login/Signup).": "Supabase Auth: email confirmation, JWT sessions, and secure login.",
        "Supabase Authentication (Secure Login/Signup)": "Supabase Auth: email confirmation, JWT sessions, and secure login",
        "Sign up, upload, analyze, and download report": "Confirm email, sign in, upload, analyze, and download a report",
        "Data layer: Supabase Auth (JWT), PostgreSQL, and Storage": "Data layer: Supabase Auth, PostgreSQL metadata, and private Storage buckets",
        "Live proof: reproducibility certificate and 76 verified numbers on the deployed site.": "Current secure sign-up screen: Supabase sends an email confirmation link before login.",
        "207": "214",
        "Supabase Auth (JWT)": "Supabase Auth (email + JWT)",
        "Supabase Auth + Storage": "Supabase Auth + PostgreSQL + Storage",
    }
    text_changes = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            for text_frame in iter_text_frames(shape):
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        updated = run.text.replace("—", ",").replace("–", "-")
                        for old, new in replacements.items():
                            updated = updated.replace(old, new)
                        if updated != run.text:
                            run.text = updated
                            text_changes += 1

    full_paragraph_replacements = {
        "• Verified end to end: sign up, upload, analyze, and download report": "• Verified account flow: sign up, confirm email, log in, upload, analyze, and download report",
        "Supabase Auth (JWT)": "Supabase Auth (email confirmation + JWT)",
        "Supabase Auth + Storage": "Supabase Auth + PostgreSQL + Storage",
    }
    for slide in prs.slides:
        for shape in slide.shapes:
            for text_frame in iter_text_frames(shape):
                for paragraph in text_frame.paragraphs:
                    if paragraph.text in full_paragraph_replacements:
                        paragraph.text = full_paragraph_replacements[paragraph.text]
                        if paragraph.text.startswith("• Verified account flow"):
                            for run in paragraph.runs:
                                run.font.size = PptPt(14)

    architecture = ARCH_OUT / "dataverse-production-architecture.png"
    slide8 = prs.slides[7]
    pictures8 = [shape for shape in slide8.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if pictures8:
        replace_ppt_picture(slide8, max(pictures8, key=lambda item: item.width * item.height), architecture)

    slide18 = prs.slides[17]
    pictures18 = [shape for shape in slide18.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if pictures18:
        replace_ppt_picture(slide18, max(pictures18, key=lambda item: item.width * item.height), UI_DIR / "latest-signup.png")

    slide14 = prs.slides[13]
    caption = slide14.shapes.add_textbox(PptInches(4.0), PptInches(6.55), PptInches(5.3), PptInches(0.28))
    paragraph = caption.text_frame.paragraphs[0]
    paragraph.text = "Table 1: DataVerse AI technology stack used in the implemented system."
    paragraph.alignment = PP_ALIGN.CENTER
    for run in paragraph.runs:
        run.font.size = PptPt(10)
        run.font.italic = True
        run.font.color.rgb = PptRGBColor(71, 85, 105)

    prs.core_properties.title = "DataVerse AI P2 Defense, Revised 2026"
    prs.core_properties.subject = "Updated Supabase architecture and current secure UI"
    PPT_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPT_OUT)
    return text_changes


def main() -> None:
    ARCH_OUT.mkdir(parents=True, exist_ok=True)
    make_architecture_diagram(ARCH_OUT / "dataverse-production-architecture.png")
    make_flow_diagram(ARCH_OUT / "dataverse-request-data-flow.png")
    changed_paragraphs, table_count = update_docx()
    ppt_text_changes = update_pptx()
    print(f"report={REPORT_OUT}")
    print(f"presentation={PPT_OUT}")
    print(f"revised_paragraphs={changed_paragraphs}")
    print(f"tables={table_count}")
    print(f"ppt_text_changes={ppt_text_changes}")


if __name__ == "__main__":
    main()
